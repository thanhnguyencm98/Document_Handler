import os
import re
import tempfile
import subprocess

# For PDF extraction
try:
    from pdfminer.high_level import extract_text as pdf_extract_text
except ImportError:
    pdf_extract_text = None

# For Office documents
try:
    import docx
    import openpyxl
    from pptx import Presentation
except ImportError:
    docx = None
    openpyxl = None
    Presentation = None

class TextExtractor:
    """
    Extracts text from various document formats.
    """
    
    def extract_text(self, file_path):
        """
        Extract text from a document file.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            str: Extracted text from the document
        """
        file_extension = os.path.splitext(file_path)[1].lower()
        
        if file_extension == ".pdf":
            return self._extract_from_pdf(file_path)
        elif file_extension == ".docx":
            return self._extract_from_docx(file_path)
        elif file_extension in [".xlsx", ".xls"]:
            return self._extract_from_excel(file_path)
        elif file_extension == ".pptx":
            return self._extract_from_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_extension}")
    
    def _extract_from_pdf(self, file_path):
        """Extract text from PDF file."""
        if pdf_extract_text:
            # Use pdfminer if available
            try:
                return pdf_extract_text(file_path)
            except Exception as e:
                print(f"Error extracting text from PDF with pdfminer: {e}")
        
        # Fallback to pdftotext if available
        try:
            # Check if pdftotext is installed
            subprocess.run(["which", "pdftotext"], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            
            # Create a temporary file for the output
            with tempfile.NamedTemporaryFile(suffix=".txt") as temp_output:
                # Run pdftotext
                subprocess.run(
                    ["pdftotext", "-layout", file_path, temp_output.name],
                    check=True,
                    stdout=subprocess.PIPE,
                    stderr=subprocess.PIPE
                )
                
                # Read the output
                with open(temp_output.name, "r") as f:
                    return f.read()
        except (subprocess.CalledProcessError, FileNotFoundError) as e:
            print(f"Error extracting text from PDF with pdftotext: {e}")
        
        # If all methods fail
        return "Error: Could not extract text from PDF. Please ensure pdfminer.six or pdftotext is installed."
    
    def _extract_from_docx(self, file_path):
        """Extract text from Word document."""
        if not docx:
            return "Error: python-docx library not installed."
        
        try:
            doc = docx.Document(file_path)
            return "\n".join(paragraph.text for paragraph in doc.paragraphs)
        except Exception as e:
            print(f"Error extracting text from Word document: {e}")
            return f"Error: Could not extract text from Word document: {str(e)}"
    
    def _extract_from_excel(self, file_path):
        """Extract text from Excel spreadsheet."""
        if not openpyxl:
            return "Error: openpyxl library not installed."
        
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            text = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text.append(f"Sheet: {sheet_name}")
                
                for row in sheet.iter_rows(values_only=True):
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    text.append("\t".join(row_values))
            
            return "\n".join(text)
        except Exception as e:
            print(f"Error extracting text from Excel file: {e}")
            return f"Error: Could not extract text from Excel file: {str(e)}"
    
    def _extract_from_pptx(self, file_path):
        """Extract text from PowerPoint presentation."""
        if not Presentation:
            return "Error: python-pptx library not installed."
        
        try:
            presentation = Presentation(file_path)
            text = []
            
            for i, slide in enumerate(presentation.slides):
                text.append(f"Slide {i+1}:")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text.append(shape.text)
            
            return "\n\n".join(text)
        except Exception as e:
            print(f"Error extracting text from PowerPoint file: {e}")
            return f"Error: Could not extract text from PowerPoint file: {str(e)}"